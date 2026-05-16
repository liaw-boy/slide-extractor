#!/usr/bin/env python3
"""Human-in-the-loop slide review UI.

Strategy: never trust a single automatic threshold. Run the extractor in
"paranoid" mode that over-samples and over-clusters, then serve a contact
sheet HTML where the user ticks the slides to keep. Rebuild the final PPTX
from the user's choices.

Why this beats pure auto-extraction for guaranteed completeness:
    Lecture videos contain genuinely ambiguous cases (animations vs new
    slides, near-identical-template consecutive slides, OCR mis-reads).
    No purely automatic threshold can resolve all of them. With a contact
    sheet, the human cost is ~30 seconds and the output is provably
    complete.
"""
from __future__ import annotations

import argparse
import http.server
import json
import logging
import socketserver
import sys
from pathlib import Path
from urllib.parse import unquote

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from slide_extractor import (
    ExtractorConfig,
    LOG,
    extract_slides,
    download_video,
)

LOG.setLevel(logging.INFO)


def build_contact_sheet(slides_dir: Path, html_path: Path) -> None:
    """Generate an HTML contact sheet with one checkbox per candidate."""
    pngs = sorted(slides_dir.glob("slide_*.png"))
    rows = []
    for i, p in enumerate(pngs, start=1):
        rel = p.name
        # Parse timestamp from filename: slide_NNN_HHhMMmSSs.png
        ts_part = rel.split("_", 2)[2].replace(".png", "")
        rows.append(
            f"""
            <label class="card">
              <input type="checkbox" data-name="{rel}" checked>
              <span class="num">#{i:03d}</span>
              <span class="ts">{ts_part}</span>
              <img src="slides/{rel}" loading="lazy" />
            </label>"""
        )
    html = f"""<!DOCTYPE html>
<html lang="zh-Hant"><head>
<meta charset="utf-8">
<title>Slide Review — {slides_dir.name}</title>
<style>
  :root {{
    --bg: #0f0f12; --card: #1c1c22; --text: #f5f5f5; --accent: #4ade80;
    --muted: #888;
  }}
  body {{
    margin: 0; padding: 0 24px 120px;
    background: var(--bg); color: var(--text);
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
  }}
  header {{
    position: sticky; top: 0; background: var(--bg); z-index: 10;
    padding: 16px 0; border-bottom: 1px solid #2a2a30;
  }}
  h1 {{ margin: 0 0 8px; font-size: 18px; font-weight: 600; }}
  .hint {{ color: var(--muted); font-size: 13px; }}
  .toolbar {{ margin-top: 10px; display: flex; gap: 8px; flex-wrap: wrap; }}
  button {{
    background: var(--card); color: var(--text); border: 1px solid #333;
    padding: 8px 14px; border-radius: 6px; cursor: pointer; font-size: 13px;
  }}
  button.primary {{ background: var(--accent); color: #000; border-color: var(--accent); }}
  .grid {{
    margin-top: 20px;
    display: grid; gap: 16px;
    grid-template-columns: repeat(auto-fill, minmax(280px, 1fr));
  }}
  .card {{
    position: relative; background: var(--card); border-radius: 10px;
    overflow: hidden; cursor: pointer; transition: transform .15s ease;
    border: 2px solid transparent;
  }}
  .card:has(input:checked) {{ border-color: var(--accent); }}
  .card:has(input:not(:checked)) {{ opacity: .35; }}
  .card img {{ width: 100%; display: block; }}
  .card input {{ position: absolute; top: 8px; right: 8px; transform: scale(1.5); }}
  .num {{
    position: absolute; top: 8px; left: 8px;
    background: rgba(0,0,0,.7); padding: 4px 8px; border-radius: 4px;
    font-size: 12px; font-weight: 600;
  }}
  .ts {{
    position: absolute; bottom: 8px; left: 8px;
    background: rgba(0,0,0,.7); padding: 4px 8px; border-radius: 4px;
    font-size: 12px; font-family: monospace;
  }}
  .status {{
    position: fixed; bottom: 20px; left: 50%; transform: translateX(-50%);
    background: var(--card); padding: 14px 20px; border-radius: 10px;
    box-shadow: 0 4px 20px rgba(0,0,0,.4);
    display: flex; gap: 12px; align-items: center;
  }}
  .count {{ color: var(--accent); font-weight: 600; }}
</style></head><body>

<header>
  <h1>Slide Review — {slides_dir.name}</h1>
  <p class="hint">勾選要保留的投影片（預設全選）。完成後按「儲存選擇」匯出最終 PPTX。</p>
  <div class="toolbar">
    <button onclick="selectAll(true)">全選</button>
    <button onclick="selectAll(false)">全取消</button>
    <button class="primary" onclick="save()">儲存選擇 → 生成 PPTX</button>
    <span class="status">已選 <span class="count" id="count">{len(pngs)}</span> / {len(pngs)} 張</span>
  </div>
</header>

<div class="grid">{"".join(rows)}</div>

<script>
const updateCount = () => {{
  const n = document.querySelectorAll('input:checked').length;
  document.getElementById('count').textContent = n;
}};
document.addEventListener('change', updateCount);
function selectAll(checked) {{
  document.querySelectorAll('input').forEach(i => i.checked = checked);
  updateCount();
}}
async function save() {{
  const kept = [...document.querySelectorAll('input:checked')].map(i => i.dataset.name);
  const r = await fetch('/finalize', {{
    method: 'POST',
    headers: {{ 'Content-Type': 'application/json' }},
    body: JSON.stringify({{ kept }})
  }});
  const result = await r.json();
  if (result.ok) {{
    alert('✓ PPTX 已生成：' + result.pptx);
  }} else {{
    alert('❌ 失敗：' + result.error);
  }}
}}
</script></body></html>"""
    html_path.write_text(html, encoding="utf-8")


def export_filtered_pptx(slides_dir: Path, kept: list[str], pptx_path: Path) -> None:
    """Build a PPTX from a subset of the candidate PNGs."""
    paths = [slides_dir / name for name in kept]
    if not paths:
        raise ValueError("no slides selected")
    with Image.open(paths[0]) as first:
        src_w, src_h = first.size
    if src_w >= src_h:
        page_w_in = 10.0
        page_h_in = 10.0 * (src_h / src_w)
    else:
        page_h_in = 10.0
        page_w_in = 10.0 * (src_w / src_h)
    prs = Presentation()
    prs.slide_width = Inches(page_w_in)
    prs.slide_height = Inches(page_h_in)
    blank = prs.slide_layouts[6]
    for p in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    prs.save(str(pptx_path))


def serve(slides_dir: Path, pptx_path: Path, port: int, bind: str) -> None:
    """HTTP server that serves the contact sheet + accepts /finalize POST."""

    class Handler(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *a, **kw):
            super().__init__(*a, directory=str(slides_dir.parent), **kw)

        def log_message(self, fmt, *args):
            print(f"[{self.address_string()}] {fmt % args}")

        def do_GET(self):
            if self.path == "/" or self.path == "/index.html":
                html = (slides_dir.parent / "_review.html").read_bytes()
                self.send_response(200)
                self.send_header("Content-Type", "text/html; charset=utf-8")
                self.send_header("Content-Length", str(len(html)))
                self.end_headers()
                self.wfile.write(html)
                return
            # serve PNGs under /slides/* by rewriting to actual dir name
            if self.path.startswith("/slides/"):
                rest = unquote(self.path[len("/slides/"):])
                self.path = "/" + slides_dir.name + "/" + rest
            super().do_GET()

        def do_POST(self):
            if self.path != "/finalize":
                self.send_error(404)
                return
            length = int(self.headers.get("Content-Length", "0"))
            body = json.loads(self.rfile.read(length))
            try:
                export_filtered_pptx(slides_dir, body.get("kept", []), pptx_path)
                resp = {"ok": True, "pptx": str(pptx_path), "n": len(body["kept"])}
            except Exception as e:  # noqa: BLE001
                resp = {"ok": False, "error": str(e)}
            payload = json.dumps(resp).encode()
            self.send_response(200 if resp["ok"] else 500)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(payload)))
            self.end_headers()
            self.wfile.write(payload)

    with socketserver.TCPServer((bind, port), Handler) as srv:
        print(f"\n▶ Review UI: http://{bind if bind != '0.0.0.0' else '<your-host>'}:{port}/")
        print("  勾選後按「儲存選擇」生成最終 PPTX。Ctrl+C 結束。\n")
        srv.serve_forever()


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(
        prog="slide-review",
        description="Human-in-the-loop slide review: over-extract then click to confirm.",
    )
    p.add_argument("source", help="Local video or YouTube URL")
    p.add_argument("-o", "--output", type=Path, default=Path.home() / "slides_output")
    p.add_argument("--cluster-jaccard", type=float, default=0.30, help="Lower = more candidates (paranoid)")
    p.add_argument("--min-duration", type=float, default=3.0, help="Lower keeps brief slides")
    p.add_argument("--sample-sec", type=float, default=2.0, help="Denser sampling")
    p.add_argument("--port", type=int, default=8901)
    p.add_argument("--bind", default="0.0.0.0")
    p.add_argument("--skip-extract", action="store_true",
                   help="Skip extraction, just serve an existing slides dir")
    args = p.parse_args(argv)

    logging.basicConfig(level=logging.INFO, format="%(message)s")

    src = args.source
    out_base: Path = args.output
    local = Path(src)
    if local.exists() and local.is_file():
        video_path = local
    else:
        video_path = download_video(src, out_base / "_video")

    title = video_path.stem
    slides_dir = out_base / title
    pptx_path = out_base / f"{title}_REVIEWED.pptx"

    if not args.skip_extract:
        if slides_dir.exists():
            import shutil
            shutil.rmtree(slides_dir)
        cfg = ExtractorConfig(
            sample_sec=args.sample_sec,
            cluster_jaccard=args.cluster_jaccard,
            min_duration_sec=args.min_duration,
        )
        paths = extract_slides(video_path, slides_dir, cfg)
        print(f"\n▶ Paranoid extraction: {len(paths)} candidates")

    build_contact_sheet(slides_dir, out_base / "_review.html")
    serve(slides_dir, pptx_path, args.port, args.bind)
    return 0


if __name__ == "__main__":
    sys.exit(main())
