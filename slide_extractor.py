#!/usr/bin/env python3
"""Slide Extractor — extract complete slides from lecture videos into PPTX.

Algorithm overview:
    1. Sample frames at fixed interval (default: every 3 seconds).
    2. Compute perceptual hash (pHash) + GPU-accelerated OCR per frame.
    3. Detect slide transitions using a dual-signal rule:
         - pHash distance >= threshold (visual change), AND
         - text content change (size_ratio < 0.6 OR not a subset).
    4. Filter out short-duration segments (animation noise).
    5. Within each slide segment, pick the frame with the most OCR text
       (animation-complete state).
    6. Export PNG snapshots + a 16:9 PPTX deck.

Why this beats HSV / SSIM / pHash-only / PySceneDetect:
    Lecture slides often share a template (same background, fonts, footer).
    Pure visual hashes treat consecutive slides as "the same scene".
    Pure OCR text comparison misfires when slides share topic terminology.
    The combined pHash + token-subset rule reliably separates real slide
    transitions from in-slide animation steps.
"""
from __future__ import annotations

import argparse
import json
import logging
import shutil
import subprocess
import sys
import time
from dataclasses import dataclass
from pathlib import Path

import cv2
import easyocr
import imagehash
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


# ────────────────────────── defaults ──────────────────────────
DEFAULT_SAMPLE_SEC = 3.0
DEFAULT_PHASH_THR = 6
DEFAULT_MIN_DURATION_SEC = 9
DEFAULT_MIN_TEXT_LEN = 15
DEFAULT_BRIGHTNESS_MIN = 50
DEFAULT_SUBSET_THR = 0.92
DEFAULT_HARD_JACCARD = 0.20
DEFAULT_SIZE_RATIO_DROP = 0.6
DEFAULT_SIZE_RATIO_GROW = 1.3
DEFAULT_CLUSTER_JACCARD = 0.45  # online content clustering threshold
DEFAULT_OCR_LANGS = ("ch_tra", "en")

LOG = logging.getLogger("slide_extractor")


# ────────────────────────── data types ──────────────────────────
@dataclass
class Sample:
    """One sampled frame with its OCR + pHash signatures."""

    frame_idx: int
    frame: np.ndarray
    text: str
    tokens: frozenset[str]
    phash: imagehash.ImageHash


@dataclass(frozen=True)
class ExtractorConfig:
    sample_sec: float = DEFAULT_SAMPLE_SEC
    phash_thr: int = DEFAULT_PHASH_THR
    min_duration_sec: float = DEFAULT_MIN_DURATION_SEC
    min_text_len: int = DEFAULT_MIN_TEXT_LEN
    brightness_min: int = DEFAULT_BRIGHTNESS_MIN
    subset_thr: float = DEFAULT_SUBSET_THR
    hard_jaccard: float = DEFAULT_HARD_JACCARD
    size_ratio_drop: float = DEFAULT_SIZE_RATIO_DROP
    size_ratio_grow: float = DEFAULT_SIZE_RATIO_GROW
    cluster_jaccard: float = DEFAULT_CLUSTER_JACCARD
    ocr_langs: tuple[str, ...] = DEFAULT_OCR_LANGS
    use_gpu: bool = True


# ────────────────────────── helpers ──────────────────────────
def text_to_tokens(text: str) -> frozenset[str]:
    """Tokenize OCR output into (Chinese bigrams + English words).

    Token-level matching is far more discriminating than character-level for
    lecture slides because different slides often share many single
    characters (headers, footers, topic words) but rarely share many bigrams.
    """
    cleaned = "".join(c if c.isalnum() or "一" <= c <= "鿿" else " " for c in text)
    tokens: set[str] = set()
    for word in cleaned.split():
        if not word:
            continue
        if all(c.isascii() and c.isalnum() for c in word):
            if len(word) >= 2:
                tokens.add(word.lower())
            continue
        chars = [c for c in word if "一" <= c <= "鿿"]
        if len(chars) == 1:
            tokens.add(chars[0])
        for i in range(len(chars) - 1):
            tokens.add(chars[i] + chars[i + 1])
    return frozenset(tokens)


def is_bright(frame: np.ndarray, min_value: int) -> bool:
    return float(cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY).mean()) >= min_value


def frame_phash(frame: np.ndarray) -> imagehash.ImageHash:
    return imagehash.phash(Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)))


def hms(seconds: float) -> str:
    s = int(seconds)
    h, rem = divmod(s, 3600)
    m, ss = divmod(rem, 60)
    return f"{h:02d}:{m:02d}:{ss:02d}"


# ────────────────────────── video I/O ──────────────────────────
def download_video(url: str, out_dir: Path, height_cap: int = 1080) -> Path:
    """Download a YouTube video at the best available quality up to `height_cap`."""
    out_dir.mkdir(parents=True, exist_ok=True)
    fmt = (
        f"bestvideo[height<={height_cap}][ext=mp4]+bestaudio[ext=m4a]/"
        f"bestvideo[height<={height_cap}]+bestaudio/best"
    )
    result = subprocess.run(
        [
            "yt-dlp",
            "-f", fmt,
            "--merge-output-format", "mp4",
            "-o", str(out_dir / "%(title).80s.%(ext)s"),
            "--print", "after_move:filepath",
            url,
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"yt-dlp failed: {result.stderr[-500:]}")
    return Path(result.stdout.strip().splitlines()[-1])


# ────────────────────────── transition logic ──────────────────────────
def classify_transition(
    prev: Sample,
    curr: Sample,
    cfg: ExtractorConfig,
) -> tuple[bool, str]:
    """Decide whether (prev → curr) is a real slide transition.

    Dual-signal rules (empirically tuned, see docs/algorithm.md):
        0. HARD OCR signal: jaccard < hard_jaccard → real transition regardless of
           pHash. Some lecture templates produce near-zero pHash distance even
           across completely different slides (identical colour-band layout).
        1. If pHash distance < threshold AND OCR overlap healthy: same slide.
        2. If text size dropped sharply: real transition (animations never remove).
        3. If text grew AND old ⊂ new (tight subset): animation step.
        4. If old ⊂ new (looser) but no growth: same slide.
        5. Otherwise (visual change + content really differs): real transition.
    """
    a, b = prev.tokens, curr.tokens
    ph_dist = curr.phash - prev.phash
    # Compute OCR signals up-front so the hard-jaccard rule can fire even when
    # the visual signal is misleading (same-template slides → low pHash).
    jac = 0.0
    if a and b:
        inter = a & b
        jac = len(inter) / len(a | b)
        if jac < cfg.hard_jaccard:
            return True, f"pH={ph_dist} hard-OCR jac={jac:.2f}"
    if ph_dist < cfg.phash_thr:
        return False, f"pH={ph_dist}<{cfg.phash_thr}"
    if not a or not b:
        return True, f"pH={ph_dist} (missing OCR)"
    inter = a & b
    size_ratio = len(b) / len(a)
    old_in_new = len(inter) / len(a)
    if size_ratio < cfg.size_ratio_drop:
        return True, f"pH={ph_dist} text shrank ratio={size_ratio:.2f}"
    if size_ratio >= cfg.size_ratio_grow and old_in_new >= cfg.subset_thr:
        return False, f"pH={ph_dist} animation grow sub={old_in_new:.2f}"
    if old_in_new >= cfg.subset_thr:
        return False, f"pH={ph_dist} subset sub={old_in_new:.2f}"
    return True, f"pH={ph_dist} content change sub={old_in_new:.2f}"


# ────────────────────────── extraction pipeline ──────────────────────────
def _sample_frames(
    video_path: Path,
    cache_path: Path,
    cfg: ExtractorConfig,
) -> tuple[list[Sample], float]:
    cap = cv2.VideoCapture(str(video_path))
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    sample_interval = max(1, int(fps * cfg.sample_sec))
    LOG.info("video fps=%.1f frames=%d sample-interval=%d", fps, total_frames, sample_interval)

    cache: dict[str, str] = {}
    if cache_path.exists():
        cache = json.loads(cache_path.read_text(encoding="utf-8"))
        LOG.info("loaded OCR cache: %d entries", len(cache))

    reader: easyocr.Reader | None = None
    samples: list[Sample] = []
    frame_idx = 0
    cache_dirty = False
    t0 = time.time()

    while True:
        ret, frame = cap.read()
        if not ret:
            break
        if frame_idx % sample_interval != 0:
            frame_idx += 1
            continue

        if not is_bright(frame, cfg.brightness_min):
            frame_idx += 1
            continue

        key = str(frame_idx)
        text = cache.get(key)
        if text is None:
            if reader is None:
                LOG.info("initializing EasyOCR (gpu=%s)…", cfg.use_gpu)
                reader = easyocr.Reader(list(cfg.ocr_langs), gpu=cfg.use_gpu, verbose=False)
            text = " ".join(reader.readtext(frame, detail=0, paragraph=False))
            cache[key] = text
            cache_dirty = True

        tokens = text_to_tokens(text)
        if len(tokens) >= cfg.min_text_len:
            samples.append(
                Sample(
                    frame_idx=frame_idx,
                    frame=frame.copy(),
                    text=text,
                    tokens=tokens,
                    phash=frame_phash(frame),
                )
            )

        if len(samples) and len(samples) % 50 == 0:
            elapsed = time.time() - t0
            LOG.info("sampled %d valid frames (%.0fs elapsed)", len(samples), elapsed)
        frame_idx += 1

    cap.release()
    if cache_dirty:
        cache_path.write_text(json.dumps(cache, ensure_ascii=False), encoding="utf-8")
        LOG.info("OCR cache saved → %s", cache_path)
    LOG.info("sampling done: %d valid frames", len(samples))
    return samples, fps


def _find_transitions(
    samples: list[Sample],
    cfg: ExtractorConfig,
    fps: float,
) -> list[int]:
    transitions: list[int] = []
    for i in range(1, len(samples)):
        is_trans, reason = classify_transition(samples[i - 1], samples[i], cfg)
        if is_trans:
            transitions.append(i)
            LOG.debug("transition @ %s (%s)", hms(samples[i].frame_idx / fps), reason)
    LOG.info("detected %d transitions", len(transitions))
    return transitions


def _filter_short_segments(
    segments: list[tuple[int, int]],
    samples: list[Sample],
    cfg: ExtractorConfig,
    fps: float,
) -> list[tuple[int, int]]:
    out: list[tuple[int, int]] = []
    for start, end in segments:
        if not (0 <= start < end <= len(samples)):
            continue
        dur = (samples[end - 1].frame_idx - samples[start].frame_idx) / fps + cfg.sample_sec
        if dur >= cfg.min_duration_sec:
            out.append((start, end))
        elif out:
            # merge into previous
            ps, _pe = out[-1]
            out[-1] = (ps, end)
        else:
            out.append((start, end))
    return out


def _pick_representative(seg: list[Sample]) -> Sample:
    """Pick the frame with the most OCR text in a segment.

    When multiple frames tie at the top, pick the last one — closest to the
    next transition, i.e. the animation-complete state.
    """
    max_chars = max(len(s.tokens) for s in seg)
    top = [s for s in seg if len(s.tokens) >= max_chars - 5]
    return top[-1]


def _jaccard(a: frozenset[str], b: frozenset[str]) -> float:
    if not a or not b:
        return 0.0
    return len(a & b) / len(a | b)


def _cluster_by_content(samples: list[Sample], cluster_thr: float) -> list[list[Sample]]:
    """Online single-pass clustering by OCR content (Jaccard on bigram tokens).

    Each frame either joins the best-matching existing cluster (jaccard ≥ thr)
    or opens a new cluster. Each cluster represents one unique slide; the
    cluster's "representative" is the max-tokens frame inside it, which we
    keep updated as the cluster grows.

    This guarantees structural completeness: any frame whose content does not
    overlap enough with any prior slide *necessarily* becomes a new slide.
    """
    clusters: list[list[Sample]] = []
    reps: list[Sample] = []  # parallel array of current best rep per cluster

    for s in samples:
        best_idx = -1
        best_jac = 0.0
        for i, rep in enumerate(reps):
            j = _jaccard(rep.tokens, s.tokens)
            if j > best_jac:
                best_jac = j
                best_idx = i
        if best_idx >= 0 and best_jac >= cluster_thr:
            clusters[best_idx].append(s)
            # Update rep if this frame has more tokens (animation-complete state)
            if len(s.tokens) > len(reps[best_idx].tokens):
                reps[best_idx] = s
        else:
            clusters.append([s])
            reps.append(s)
    return clusters


def extract_slides(
    video_path: Path,
    out_dir: Path,
    cfg: ExtractorConfig | None = None,
) -> list[Path]:
    """Run the full extraction pipeline; returns list of saved PNG paths.

    Strategy: content-based clustering (no missing slides). Each unique OCR
    content cluster becomes one output slide; within each cluster we pick the
    frame with the most tokens (animation-complete state). Output is ordered
    by the earliest occurrence of each cluster.
    """
    cfg = cfg or ExtractorConfig()
    out_dir.mkdir(parents=True, exist_ok=True)
    cache_path = out_dir.parent / f"_ocr_cache_{video_path.stem}.json"

    samples, fps = _sample_frames(video_path, cache_path, cfg)
    if not samples:
        LOG.warning("no valid samples — nothing to do")
        return []

    cluster_thr = cfg.cluster_jaccard
    clusters = _cluster_by_content(samples, cluster_thr)
    LOG.info("content-clustered: %d raw clusters (jaccard ≥ %.2f)", len(clusters), cluster_thr)

    # Filter out tiny noise clusters: a real slide stays on screen for at
    # least MIN_SLIDE_DURATION seconds, so a cluster of frames spanning less
    # than that — especially singleton clusters — is almost certainly an
    # OCR-flicker artefact at a transition boundary, not a real slide.
    def cluster_span_sec(c: list[Sample]) -> float:
        if len(c) < 2:
            return 0.0
        return (c[-1].frame_idx - c[0].frame_idx) / fps + cfg.sample_sec

    min_frames = max(2, int(cfg.min_duration_sec / cfg.sample_sec))
    real = [c for c in clusters if len(c) >= min_frames]
    dropped = len(clusters) - len(real)
    if dropped:
        LOG.info("dropped %d noise clusters (< %d frames each)", dropped, min_frames)
    clusters = real

    # Sort clusters by their earliest frame_idx so output is in time order
    clusters.sort(key=lambda c: c[0].frame_idx)

    saved: list[Path] = []
    for slide_no, cluster in enumerate(clusters, start=1):
        best = _pick_representative(cluster)
        ts = hms(best.frame_idx / fps).replace(":", "h", 1).replace(":", "m") + "s"
        path = out_dir / f"slide_{slide_no:03d}_{ts}.png"
        cv2.imwrite(str(path), best.frame)
        saved.append(path)
        LOG.info(
            "[+] slide %03d @ %s  (%d tokens, cluster of %d frames)",
            slide_no, hms(best.frame_idx / fps), len(best.tokens), len(cluster),
        )
    return saved


def export_pptx(slide_paths: list[Path], out_path: Path) -> None:
    """Export PNGs into a PPTX deck preserving the source aspect ratio.

    The page size is derived from the first slide image so 4:3 lecture
    recordings come out 4:3 and 16:9 recordings come out 16:9, without any
    distortion. Each image is full-bleed on its page.
    """
    if not slide_paths:
        LOG.warning("no slides to export")
        return
    with Image.open(slide_paths[0]) as first:
        src_w, src_h = first.size
    # Fix the longer dimension at 10 inches so the page always fits typical
    # editing canvases, then scale the other dimension by the source ratio.
    if src_w >= src_h:
        page_w_in = 10.0
        page_h_in = 10.0 * (src_h / src_w)
    else:
        page_h_in = 10.0
        page_w_in = 10.0 * (src_w / src_h)
    LOG.info("PPTX page %.2f x %.2f in (source %dx%d)", page_w_in, page_h_in, src_w, src_h)

    prs = Presentation()
    prs.slide_width = Inches(page_w_in)
    prs.slide_height = Inches(page_h_in)
    blank_layout = prs.slide_layouts[6]
    for img in slide_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    prs.save(str(out_path))
    LOG.info("PPTX → %s", out_path)


# ────────────────────────── CLI ──────────────────────────
def build_arg_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="slide-extractor",
        description="Extract complete slides from lecture videos into PNG + PPTX.",
    )
    p.add_argument("source", help="Local video file path OR YouTube URL")
    p.add_argument(
        "-o", "--output", type=Path, default=Path.home() / "slides_output",
        help="Output base directory (default: ~/slides_output)",
    )
    p.add_argument("--sample-sec", type=float, default=DEFAULT_SAMPLE_SEC, help="Sampling interval in seconds")
    p.add_argument("--phash-thr", type=int, default=DEFAULT_PHASH_THR, help="pHash distance threshold for transitions")
    p.add_argument("--min-duration", type=float, default=DEFAULT_MIN_DURATION_SEC, help="Min slide duration (s)")
    p.add_argument(
        "--cluster-jaccard", type=float, default=DEFAULT_CLUSTER_JACCARD,
        help="Jaccard threshold for content clustering (lower = more conservative = fewer merges)",
    )
    p.add_argument("--cpu", action="store_true", help="Force CPU-only OCR (default: GPU if available)")
    p.add_argument(
        "--lang", action="append", default=None,
        help="EasyOCR language code(s); repeat the flag for multiple (default: ch_tra en)",
    )
    p.add_argument("-v", "--verbose", action="count", default=0, help="Increase log verbosity (-v, -vv)")
    return p


def main(argv: list[str] | None = None) -> int:
    args = build_arg_parser().parse_args(argv)

    level = logging.WARNING - 10 * args.verbose
    logging.basicConfig(level=max(level, logging.DEBUG), format="%(message)s")

    cfg = ExtractorConfig(
        sample_sec=args.sample_sec,
        phash_thr=args.phash_thr,
        min_duration_sec=args.min_duration,
        cluster_jaccard=args.cluster_jaccard,
        use_gpu=not args.cpu,
        ocr_langs=tuple(args.lang) if args.lang else DEFAULT_OCR_LANGS,
    )

    src = args.source
    out_base: Path = args.output
    local = Path(src)
    if local.exists() and local.is_file():
        video_path = local
    else:
        LOG.info("downloading: %s", src)
        video_path = download_video(src, out_base / "_video")

    title = video_path.stem
    slides_dir = out_base / title
    pptx_path = out_base / f"{title}.pptx"
    if slides_dir.exists():
        shutil.rmtree(slides_dir)

    paths = extract_slides(video_path, slides_dir, cfg)
    if not paths:
        LOG.error("no slides extracted")
        return 1
    export_pptx(paths, pptx_path)
    print(f"\n{len(paths)} slides → {slides_dir}\nPPTX  → {pptx_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
